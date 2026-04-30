"""
NeuroClaw Web UI Server

Serves a browser-based chat interface at http://localhost:7080 by default.

Usage
-----
    # Preferred: via the main agent entry point
    python core/agent/main.py --web [--port 7080] [--host 127.0.0.1]

    # Or run the web server directly
    python core/web/server.py [--port 7080] [--host 127.0.0.1]

Dependencies (install via the setup wizard or manually)
---------------------------------------------------------
    pip install "fastapi[standard]" uvicorn
"""
from __future__ import annotations

import argparse
import asyncio
import importlib.util
import io
import json
import os
import subprocess
import time
import re
import queue as stdlib_queue
import sys
import threading
from pathlib import Path
from typing import Any

REPO_ROOT = Path(__file__).parent.parent.parent.resolve()
STATIC_DIR = Path(__file__).parent / "static"
SHELL_STATUS_FILE = Path("/tmp/neuroclaw_claw_shell_status.json")
AGENT_SHELL_STATUS_FILE = Path("/tmp/neuroclaw_agent_shell_status.json")

DEFAULT_HOST = "0.0.0.0"
DEFAULT_PORT = 7080
ATTACHMENT_MAX_FILE_BYTES = 15 * 1024 * 1024
ATTACHMENT_MAX_EMBED_CHARS = 12000
SUPPORTED_ATTACHMENT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".json", ".yaml", ".yml", ".csv", ".tsv",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".sh", ".bash", ".zsh", ".sql",
    ".html", ".css", ".xml", ".log", ".rst", ".ini", ".toml", ".cfg",
    ".pdf", ".docx", ".xlsx", ".pptx",
}


def _fallback_title_from_user_text(text: str) -> str:
    """Create a short deterministic title when LLM title generation fails."""
    cleaned = re.sub(r"\s+", " ", str(text or "")).strip()
    if not cleaned:
        return "New Chat"
    words = cleaned.split(" ")[:8]
    title = " ".join(words)
    return title[:64] if title else "New Chat"


def _sanitize_title(raw: str, user_text: str) -> str:
    """Normalize model output into a single short title line."""
    t = str(raw or "").strip()
    if not t:
        return _fallback_title_from_user_text(user_text)

    t = t.splitlines()[0].strip()
    t = t.strip("\"'` ")
    t = re.sub(r"^[\-\*\d\.)\s]+", "", t)
    t = re.sub(r"\s+", " ", t).strip(" .:-")
    if not t:
        return _fallback_title_from_user_text(user_text)
    return t[:64]


def _safe_skill_summary_fallback(skill_name: str, description: str) -> dict[str, str]:
    """Return conservative bilingual fallback summary when LLM summarization fails."""
    base_en = (description or "").strip()
    if not base_en:
        base_en = f"{skill_name} provides a specialized workflow for task execution in NeuroClaw."
    base_en = re.sub(r"\s+", " ", base_en).strip()
    if len(base_en) > 220:
        base_en = base_en[:217].rstrip() + "..."
    base_zh = f"该技能围绕「{skill_name}」提供专用流程支持，可用于相关任务的执行与组织。"
    return {"en": base_en, "zh": base_zh}


def _parse_bilingual_summary_json(raw: str) -> dict[str, str] | None:
    """Parse model output into {'en':..., 'zh':...} with defensive handling."""
    text = str(raw or "").strip()
    if not text:
        return None

    try:
        obj = json.loads(text)
        en = str(obj.get("en", "")).strip()
        zh = str(obj.get("zh", "")).strip()
        if en and zh:
            return {"en": en, "zh": zh}
    except Exception:
        pass

    match = re.search(r"\{[\s\S]*\}", text)
    if match:
        try:
            obj = json.loads(match.group(0))
            en = str(obj.get("en", "")).strip()
            zh = str(obj.get("zh", "")).strip()
            if en and zh:
                return {"en": en, "zh": zh}
        except Exception:
            return None
    return None


def _strip_frontmatter(text: str) -> str:
    """Remove optional YAML front-matter from SKILL.md text."""
    return re.sub(r"^---\s*\n.*?\n---\s*\n", "", text, flags=re.DOTALL)


def _read_shell_status() -> dict[str, Any] | None:
    """Read transient shell status from agent-shell first, then tmux claw-shell."""
    if AGENT_SHELL_STATUS_FILE.exists():
        try:
            data = json.loads(AGENT_SHELL_STATUS_FILE.read_text(encoding="utf-8"))
            command = str(data.get("command", "")).strip()
            started_at = data.get("started_at")
            pid = int(data.get("pid", 0))
            if command and isinstance(started_at, (int, float)) and pid > 0:
                alive = False
                try:
                    os.kill(pid, 0)
                    alive = True
                except Exception:
                    alive = False

                if alive:
                    return {
                        "active": True,
                        "source": "agent_shell",
                        "command": command,
                        "current_command": "agent_shell",
                        "started_at": int(started_at),
                        "elapsed_ms": max(0, int((time.time() * 1000) - int(started_at))),
                    }
        except Exception:
            pass
        try:
            AGENT_SHELL_STATUS_FILE.unlink()
        except Exception:
            pass

    if not SHELL_STATUS_FILE.exists():
        return None

    try:
        data = json.loads(SHELL_STATUS_FILE.read_text(encoding="utf-8"))
    except Exception:
        try:
            SHELL_STATUS_FILE.unlink()
        except Exception:
            pass
        return None

    command = str(data.get("command", "")).strip()
    started_at = data.get("started_at")
    if not command or not isinstance(started_at, (int, float)):
        try:
            SHELL_STATUS_FILE.unlink()
        except Exception:
            pass
        return None

    try:
        proc = subprocess.run(
            ["tmux", "display-message", "-p", "-t", "claw", "#{pane_current_command}"],
            capture_output=True,
            text=True,
            check=False,
            timeout=2,
        )
        current = (proc.stdout or "").strip().lower()
    except Exception:
        current = ""

    shell_names = {"bash", "zsh", "fish", "sh", "dash", "tmux"}
    active = bool(current) and current not in shell_names
    elapsed_ms = max(0, int((time.time() * 1000) - int(started_at)))

    if not active:
        try:
            SHELL_STATUS_FILE.unlink()
        except Exception:
            pass
        return None

    return {
        "active": True,
        "command": command,
        "current_command": current,
        "started_at": int(started_at),
        "elapsed_ms": elapsed_ms,
    }


def _skill_md_excerpt(skill_md: Path, max_chars: int = 1800) -> str:
    """Return a compact excerpt from SKILL.md for prompt context injection."""
    try:
        raw = skill_md.read_text(encoding="utf-8")
    except Exception:
        return ""
    body = _strip_frontmatter(raw)
    body = re.sub(r"```[\s\S]*?```", "", body)
    body = re.sub(r"\n{3,}", "\n\n", body).strip()
    if len(body) > max_chars:
        body = body[:max_chars].rstrip() + "\n..."
    return body


def _selected_skills_context(selected_names: list[str], skills: list[dict[str, Any]]) -> str:
    """Build selected-skill references from discovered SKILL.md files."""
    if not selected_names:
        return ""

    by_name: dict[str, dict[str, Any]] = {}
    for s in skills:
        n = str(s.get("name", "")).strip().lower()
        if n:
            by_name[n] = s

    chunks: list[str] = []
    for raw_name in selected_names:
        key = raw_name.strip().lower()
        if not key:
            continue
        s = by_name.get(key)
        if not s:
            continue
        excerpt = _skill_md_excerpt(Path(s.get("skill_md")))
        if not excerpt:
            continue
        chunks.append(
            f"[Skill: {s.get('name', raw_name)}]\n"
            f"Description: {s.get('description', '')}\n"
            f"SKILL.md excerpt:\n{excerpt}"
        )

    return "\n\n".join(chunks)


def _attachment_extension(filename: str) -> str:
    return Path(str(filename or "")).suffix.lower().strip()


def _truncate_text(text: str, max_chars: int = ATTACHMENT_MAX_EMBED_CHARS) -> tuple[str, bool]:
    raw = str(text or "")
    if len(raw) <= max_chars:
        return raw, False
    return raw[:max_chars].rstrip() + "\n[Content truncated]", True


def _decode_text_attachment(raw: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return raw.decode(encoding)
        except Exception:
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_pdf_text(raw: bytes) -> str:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(io.BytesIO(raw))
    chunks: list[str] = []
    for page in reader.pages[:120]:
        txt = (page.extract_text() or "").strip()
        if txt:
            chunks.append(txt)
    return "\n\n".join(chunks)


def _extract_docx_text(raw: bytes) -> str:
    from docx import Document  # type: ignore

    doc = Document(io.BytesIO(raw))
    chunks: list[str] = []

    for para in doc.paragraphs:
        txt = (para.text or "").strip()
        if txt:
            chunks.append(txt)

    for table in doc.tables:
        for row in table.rows:
            vals = [str(cell.text or "").strip() for cell in row.cells]
            vals = [v for v in vals if v]
            if vals:
                chunks.append(" | ".join(vals))

    return "\n".join(chunks)


def _extract_xlsx_text(raw: bytes) -> str:
    from openpyxl import load_workbook  # type: ignore

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets[:12]:
        lines.append(f"[Sheet] {ws.title}")
        row_count = 0
        for row in ws.iter_rows(max_row=200, max_col=30, values_only=True):
            cells = ["" if v is None else str(v).strip() for v in row]
            if not any(cells):
                continue
            lines.append("\t".join(cells))
            row_count += 1
            if row_count >= 200:
                lines.append("[Rows truncated]")
                break
        lines.append("")
    return "\n".join(lines).strip()


def _extract_pptx_text(raw: bytes) -> str:
    from pptx import Presentation  # type: ignore

    prs = Presentation(io.BytesIO(raw))
    lines: list[str] = []
    for idx, slide in enumerate(prs.slides[:80], start=1):
        lines.append(f"[Slide {idx}]")
        for shape in slide.shapes:
            txt = ""
            try:
                txt = (shape.text or "").strip()  # type: ignore[attr-defined]
            except Exception:
                txt = ""
            if txt:
                lines.append(txt)
        lines.append("")
    return "\n".join(lines).strip()


def _parse_attachment_content(filename: str, content_type: str, raw: bytes) -> dict[str, Any]:
    ext = _attachment_extension(filename)
    if ext not in SUPPORTED_ATTACHMENT_EXTENSIONS:
        return {
            "ok": False,
            "error": (
                f"Unsupported extension '{ext or 'unknown'}'. "
                "Supported: " + ", ".join(sorted(SUPPORTED_ATTACHMENT_EXTENSIONS))
            ),
        }

    if ext in {
        ".txt", ".md", ".markdown", ".json", ".yaml", ".yml", ".csv", ".tsv",
        ".py", ".js", ".ts", ".tsx", ".jsx", ".sh", ".bash", ".zsh", ".sql",
        ".html", ".css", ".xml", ".log", ".rst", ".ini", ".toml", ".cfg",
    }:
        text = _decode_text_attachment(raw)
    elif ext == ".pdf":
        try:
            text = _extract_pdf_text(raw)
        except ImportError:
            return {"ok": False, "error": "PDF parser missing. Install: pypdf"}
    elif ext == ".docx":
        try:
            text = _extract_docx_text(raw)
        except ImportError:
            return {"ok": False, "error": "DOCX parser missing. Install: python-docx"}
    elif ext == ".xlsx":
        try:
            text = _extract_xlsx_text(raw)
        except ImportError:
            return {"ok": False, "error": "XLSX parser missing. Install: openpyxl"}
    elif ext == ".pptx":
        try:
            text = _extract_pptx_text(raw)
        except ImportError:
            return {"ok": False, "error": "PPTX parser missing. Install: python-pptx"}
    else:
        return {"ok": False, "error": "Unsupported file type"}

    cleaned = str(text or "").strip()
    if not cleaned:
        return {"ok": False, "error": "No readable text found in file"}

    truncated_text, truncated = _truncate_text(cleaned)
    return {
        "ok": True,
        "text": truncated_text,
        "truncated": truncated,
        "ext": ext,
        "content_type": content_type,
    }


def _normalize_skill_token(text: str) -> str:
    """Normalize a skill token for robust mention matching."""
    t = str(text or "").lower()
    t = re.sub(r"[^a-z0-9]+", "", t)
    return t


def _infer_skills_from_user_text(user_text: str, skills: list[dict[str, Any]]) -> list[str]:
    """
    Infer referenced skills from free-form user text.

    Matches both skill name and skill directory name in normalized form.
    """
    norm_msg = _normalize_skill_token(user_text)
    if not norm_msg:
        return []

    inferred: list[str] = []
    for s in skills:
        skill_name = str(s.get("name", "")).strip()
        if not skill_name:
            continue
        dir_name = ""
        try:
            dir_name = Path(s.get("path")).name
        except Exception:
            dir_name = ""

        candidates = [skill_name, dir_name]
        for c in candidates:
            token = _normalize_skill_token(c)
            if token and token in norm_msg:
                inferred.append(skill_name)
                break

    # Preserve order and remove duplicates
    seen: set[str] = set()
    out: list[str] = []
    for name in inferred:
        key = name.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(name)
    return out[:5]


# ── Module import helpers ──────────────────────────────────────────────────────

def _import_from_path(module_name: str, path: Path) -> Any:
    """Import a Python module from an absolute file path (handles hyphenated dirs)."""
    spec = importlib.util.spec_from_file_location(module_name, path)
    if spec is None or spec.loader is None:
        raise ImportError(f"Cannot load module from {path}")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)  # type: ignore[union-attr]
    return mod


# ── Dependency check ───────────────────────────────────────────────────────────

def _require_webdeps() -> None:
    """Raise a descriptive RuntimeError if fastapi/uvicorn are not installed."""
    missing = []
    for pkg in ("fastapi", "uvicorn"):
        try:
            __import__(pkg)
        except ImportError:
            missing.append(pkg)
    if missing:
        raise RuntimeError(
            f"Web UI dependencies not installed: {', '.join(missing)}\n"
            "Install with:  pip install 'fastapi[standard]' uvicorn\n"
            "Or re-run the installer:  python installer/setup.py"
        )


# ── Streaming helpers ──────────────────────────────────────────────────────────

async def _stream_openai(
    websocket: Any, llm_client: Any, model: str, history: list[dict]
) -> str:
    """
    Stream an OpenAI response chunk-by-chunk.

    Uses a producer thread + asyncio queue so the blocking OpenAI iterator
    does not stall the event loop while still delivering incremental updates
    to the browser.
    """
    q: stdlib_queue.Queue[tuple[str, str | None]] = stdlib_queue.Queue()

    def _produce() -> None:
        try:
            stream = llm_client.chat.completions.create(
                model=model, messages=history, stream=True
            )
            for chunk in stream:
                content = ""
                if chunk.choices and chunk.choices[0].delta:
                    content = chunk.choices[0].delta.content or ""
                if content:
                    q.put(("chunk", content))
        except Exception as exc:
            q.put(("error", str(exc)))
        finally:
            q.put(("done", None))

    threading.Thread(target=_produce, daemon=True).start()

    full = ""
    while True:
        try:
            kind, data = q.get_nowait()
        except stdlib_queue.Empty:
            await asyncio.sleep(0.01)
            continue

        if kind == "chunk":
            full += data  # type: ignore[operator]
            await websocket.send_text(json.dumps({"type": "chunk", "content": data}))
        elif kind == "error":
            raise RuntimeError(data)
        else:  # "done"
            break

    return full


async def _respond(websocket: Any, session: Any) -> str:
    """
    Generate a reply for the latest message in session.history.

    Streams chunks to the browser for OpenAI; falls back to asyncio.to_thread
    for other backends (Anthropic, local).  Always sends a final
    ``{"type": "done", "content": "…"}`` frame.
    """
    provider = session.env.get("llm_backend", {}).get("provider", "openai")
    model = session.env.get("llm_backend", {}).get("model", "gpt-4o")

    if provider == "openai" and session._llm is not None:
        full = await _stream_openai(websocket, session._llm, model, session.history)
    else:
        # Non-streaming fallback: run blocking _chat() in a thread pool
        full = await asyncio.to_thread(session._chat)

    await websocket.send_text(json.dumps({"type": "done", "content": full}))
    return full


# ── FastAPI application factory ────────────────────────────────────────────────

def create_app() -> Any:
    """Build and return the FastAPI application object."""
    _require_webdeps()

    from fastapi import FastAPI, File, UploadFile, WebSocket, WebSocketDisconnect  # type: ignore
    from fastapi.responses import FileResponse, HTMLResponse, JSONResponse  # type: ignore
    from fastapi.staticfiles import StaticFiles  # type: ignore

    # Ensure repo root is on sys.path so `from core.agent.main import …` resolves
    if str(REPO_ROOT) not in sys.path:
        sys.path.insert(0, str(REPO_ROOT))

    from core.agent.main import (  # type: ignore[import]
        AgentSession,
        build_llm_client,
        load_environment,
        save_environment,
    )

    # SkillLoader lives in a hyphenated directory (core/skill-loader/), so we
    # cannot use a regular package import — use importlib instead.
    _loader_mod = _import_from_path(
        "neuroclaw_skill_loader",
        REPO_ROOT / "core" / "skill-loader" / "loader.py",
    )
    SkillLoader = _loader_mod.SkillLoader

    app = FastAPI(title="NeuroClaw Web UI", docs_url=None, redoc_url=None)

    # ── Static files ────────────────────────────────────────────────────────────
    if STATIC_DIR.exists():
        app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")
    materials_dir = REPO_ROOT / "materials"
    if materials_dir.exists():
        app.mount("/materials", StaticFiles(directory=str(materials_dir)), name="materials")

    # ── HTTP endpoints ──────────────────────────────────────────────────────────

    @app.get("/")
    async def root() -> Any:
        index = STATIC_DIR / "index.html"
        if index.exists():
            return FileResponse(str(index))
        return HTMLResponse(
            "<h1>NeuroClaw Web UI</h1>"
            f"<p>Static files not found at <code>{STATIC_DIR}</code>.</p>",
            status_code=500,
        )

    @app.get("/api/health")
    async def health() -> dict:
        return {"status": "ok"}

    @app.get("/api/shell/status")
    async def shell_status() -> dict:
        status = _read_shell_status()
        if not status:
            return {"active": False}
        return status

    @app.get("/api/skills")
    async def list_skills() -> dict:
        loader = SkillLoader(REPO_ROOT / "skills")
        skills = loader.load_all()
        return {
            "skills": [
                {
                    "name": s["name"],
                    "description": s.get("description", ""),
                    "summary_en": s.get("summary_en", ""),
                    "summary_zh": s.get("summary_zh", ""),
                }
                for s in skills
            ]
        }

    @app.post("/api/attachments/parse")
    async def parse_attachments(files: list[UploadFile] = File(...)) -> Any:
        if not files:
            return JSONResponse({"type": "error", "message": "No files uploaded"}, status_code=400)

        parsed_files: list[dict[str, Any]] = []
        for upload in files[:24]:
            name = str(upload.filename or "untitled")
            content_type = str(upload.content_type or "unknown")
            raw = await upload.read()
            size = len(raw)

            item: dict[str, Any] = {
                "name": name,
                "size": size,
                "content_type": content_type,
            }

            if size > ATTACHMENT_MAX_FILE_BYTES:
                item.update(
                    {
                        "ok": False,
                        "error": (
                            f"File too large ({size} bytes). "
                            f"Limit is {ATTACHMENT_MAX_FILE_BYTES} bytes."
                        ),
                    }
                )
                parsed_files.append(item)
                continue

            try:
                item.update(_parse_attachment_content(name, content_type, raw))
            except Exception as exc:
                item.update({"ok": False, "error": f"Parse failed: {exc}"})
            parsed_files.append(item)

        return {
            "type": "done",
            "files": parsed_files,
            "max_file_bytes": ATTACHMENT_MAX_FILE_BYTES,
            "supported_extensions": sorted(SUPPORTED_ATTACHMENT_EXTENSIONS),
        }

    @app.get("/api/env")
    async def get_env() -> dict:
        """Return non-sensitive parts of the runtime environment config."""
        env = load_environment()
        llm = env.get("llm_backend", {})
        provider = llm.get("provider", "unknown")
        api_key_env = llm.get("api_key_env", "")
        api_key_present = bool(api_key_env and __import__('os').environ.get(api_key_env))
        return {
            "provider": provider,
            "model": llm.get("model", "unknown"),
            "available_models": llm.get("available_models", []),
            "cuda_device": env.get("cuda", {}).get("device", "cpu"),
            "setup_type": env.get("setup_type", "unknown"),
            "conda_env": env.get("conda_env"),
            "api_key_present": api_key_present,
        }

    @app.post("/api/env/model")
    async def set_model(payload: dict) -> Any:
        """Switch current provider/model to one of the configured options."""
        provider = str(payload.get("provider", "")).strip()
        model = str(payload.get("model", "")).strip()
        if not provider or not model:
            return JSONResponse(
                {"type": "error", "message": "provider and model are required"},
                status_code=400,
            )

        env = load_environment()
        llm = env.setdefault("llm_backend", {})
        available_models = llm.get("available_models", [])
        if not any(
            isinstance(item, dict)
            and str(item.get("provider", "")).strip() == provider
            and str(item.get("model", item.get("id", item.get("name", "")))).strip() == model
            for item in available_models
        ):
            return JSONResponse(
                {"type": "error", "message": "Requested provider/model is not configured"},
                status_code=400,
            )

        llm["provider"] = provider
        llm["model"] = model
        save_environment(env)
        return {
            "type": "done",
            "provider": provider,
            "model": model,
            "available_models": llm.get("available_models", []),
        }

    @app.post("/api/chat")
    async def chat_http(payload: dict) -> Any:
        """HTTP fallback for chat when WebSocket is unavailable."""
        user_text = str(payload.get("message", "")).strip()
        raw_history = payload.get("history", [])
        raw_selected_skills = payload.get("selected_skills", [])
        if not user_text:
            return JSONResponse({"type": "error", "message": "Empty message"}, status_code=400)

        session = AgentSession()

        try:
            loader = SkillLoader(REPO_ROOT / "skills")
            skills = loader.load_all()
        except Exception:
            skills = []

        selected_skills: list[str] = []
        if isinstance(raw_selected_skills, list):
            selected_skills = [
                str(x).strip() for x in raw_selected_skills if str(x).strip()
            ]
        if not selected_skills:
            selected_skills = _infer_skills_from_user_text(user_text, skills)

        env_file = REPO_ROOT / "neuroclaw_environment.json"
        if not env_file.exists():
            return JSONResponse({
                "type": "error",
                "message": "neuroclaw_environment.json not found. Run python installer/setup.py to configure NeuroClaw.",
            }, status_code=400)

        try:
            session.set_llm_client(build_llm_client(session.env))
        except Exception as exc:
            return JSONResponse({"type": "error", "message": f"LLM backend error: {exc}"}, status_code=500)

        soul_path = REPO_ROOT / "SOUL.md"
        soul = soul_path.read_text(encoding="utf-8") if soul_path.exists() else ""
        skill_names = ", ".join(s["name"] for s in skills)
        session.history = [{"role": "system", "content": f"{soul}\n\nLoaded skills: {skill_names}"}]

        if isinstance(raw_history, list):
            for msg in raw_history:
                if not isinstance(msg, dict):
                    continue
                role = str(msg.get("role", "")).strip().lower()
                content = str(msg.get("content", "")).strip()
                if role not in {"user", "assistant"} or not content:
                    continue
                session.history.append({"role": role, "content": content})

        selected_ctx = _selected_skills_context(selected_skills, skills)
        user_payload = user_text
        if selected_ctx:
            user_payload = (
                f"{user_text}\n\n"
                "[Selected skill references from local SKILL.md files]\n"
                f"{selected_ctx}"
            )

        session.history.append({"role": "user", "content": user_payload})

        reply = await asyncio.to_thread(session._chat)
        llm_cfg = session.env.get("llm_backend", {})
        provider_used = str(llm_cfg.get("provider", "unknown"))
        model_used = str(llm_cfg.get("model", "unknown"))
        return {
            "type": "done",
            "content": reply,
            "provider_used": provider_used,
            "model_used": model_used,
        }

    @app.post("/api/chat/title")
    async def chat_title(payload: dict) -> Any:
        """Generate a concise chat title using the configured default model."""
        user_text = str(payload.get("user", "")).strip()
        assistant_text = str(payload.get("assistant", "")).strip()
        if not user_text and not assistant_text:
            return JSONResponse({"type": "error", "message": "Missing conversation content"}, status_code=400)

        env_file = REPO_ROOT / "neuroclaw_environment.json"
        if not env_file.exists():
            return JSONResponse({"type": "error", "message": "Environment not configured"}, status_code=400)

        session = AgentSession()
        try:
            session.set_llm_client(build_llm_client(session.env))
        except Exception as exc:
            return JSONResponse({"type": "error", "message": f"LLM backend error: {exc}"}, status_code=500)

        title_system = (
            "You are a conversation title generator. "
            "Return exactly one short title in plain text, 3-10 words, no quotes, no markdown."
        )
        convo = f"User message:\n{user_text}\n\nAssistant reply:\n{assistant_text}"
        session.history = [
            {"role": "system", "content": title_system},
            {"role": "user", "content": convo},
        ]

        try:
            raw_title = await asyncio.to_thread(session._chat)
            title = _sanitize_title(raw_title, user_text)
        except Exception:
            title = _fallback_title_from_user_text(user_text)

        return {"type": "done", "title": title}

    @app.post("/api/skills/summary")
    async def skill_summary(payload: dict) -> Any:
        """Summarize one SKILL.md into bilingual 1-5 sentence summaries."""
        skill_name = str(payload.get("name", "")).strip()
        if not skill_name:
            return JSONResponse({"type": "error", "message": "Missing skill name"}, status_code=400)

        try:
            loader = SkillLoader(REPO_ROOT / "skills")
            skills = loader.load_all()
        except Exception:
            skills = []

        target = None
        for s in skills:
            if str(s.get("name", "")).strip().lower() == skill_name.lower():
                target = s
                break
        if target is None:
            return JSONResponse({"type": "error", "message": f"Skill not found: {skill_name}"}, status_code=404)

        skill_md = Path(target.get("skill_md"))
        if not skill_md.exists():
            fb = _safe_skill_summary_fallback(skill_name, str(target.get("description", "")))
            return {"type": "done", "summary_en": fb["en"], "summary_zh": fb["zh"]}

        md_text = skill_md.read_text(encoding="utf-8")
        excerpt = md_text[:7000]

        session = AgentSession()
        try:
            session.set_llm_client(build_llm_client(session.env))
        except Exception:
            fb = _safe_skill_summary_fallback(skill_name, str(target.get("description", "")))
            return {"type": "done", "summary_en": fb["en"], "summary_zh": fb["zh"]}

        prompt = (
            "Summarize the following SKILL.md into concise bilingual summaries.\n"
            "Requirements:\n"
            "1) Return valid JSON only, no markdown, no extra text.\n"
            "2) JSON keys: en, zh.\n"
            "3) en: 1-5 complete English sentences.\n"
            "4) zh: 1-5 complete Chinese sentences.\n"
            "5) Focus on what this skill can do and when to use it.\n\n"
            f"Skill Name: {skill_name}\n"
            f"SKILL.md Content:\n{excerpt}"
        )
        session.history = [
            {
                "role": "system",
                "content": "You are a precise technical summarizer. Output strict JSON only.",
            },
            {"role": "user", "content": prompt},
        ]

        try:
            raw = await asyncio.to_thread(session._chat)
            parsed = _parse_bilingual_summary_json(raw)
            if not parsed:
                raise ValueError("Invalid summary JSON")
            return {
                "type": "done",
                "summary_en": parsed["en"],
                "summary_zh": parsed["zh"],
            }
        except Exception:
            fb = _safe_skill_summary_fallback(skill_name, str(target.get("description", "")))
            return {"type": "done", "summary_en": fb["en"], "summary_zh": fb["zh"]}

    # ── WebSocket chat endpoint ─────────────────────────────────────────────────

    @app.websocket("/ws/chat")
    async def chat_endpoint(websocket: WebSocket) -> None:
        try:
            await websocket.accept()
            print("[WS] Client connected", flush=True)

            # Create a per-connection agent session
            session = AgentSession()

            # Load skills
            try:
                loader = SkillLoader(REPO_ROOT / "skills")
                skills = loader.load_all()
            except Exception:
                skills = []

            # Send init metadata
            llm_cfg = session.env.get("llm_backend", {})
            await websocket.send_text(json.dumps({
                "type": "init",
                "skills": [
                    {
                        "name": s["name"],
                        "description": s.get("description", ""),
                        "summary_en": s.get("summary_en", ""),
                        "summary_zh": s.get("summary_zh", ""),
                    }
                    for s in skills
                ],
                "provider": llm_cfg.get("provider", "unconfigured"),
                "model": llm_cfg.get("model", "unconfigured"),
            }))

            env_file = REPO_ROOT / "neuroclaw_environment.json"
            if not env_file.exists():
                await websocket.send_text(json.dumps({
                    "type": "error",
                    "message": (
                        "neuroclaw_environment.json not found. "
                        "Run python installer/setup.py to configure NeuroClaw."
                    ),
                }))
                await websocket.close()
                return

            # Initialise LLM client
            try:
                session.set_llm_client(build_llm_client(session.env))
            except Exception as exc:
                await websocket.send_text(json.dumps({
                    "type": "error",
                    "message": f"LLM backend error: {exc}",
                }))

            # Build system prompt
            soul_path = REPO_ROOT / "SOUL.md"
            soul = soul_path.read_text(encoding="utf-8") if soul_path.exists() else ""
            skill_names = ", ".join(s["name"] for s in skills)
            session.history = [
                {"role": "system", "content": f"{soul}\n\nLoaded skills: {skill_names}"}
            ]

            # Main chat loop
            try:
                while True:
                    raw = await websocket.receive_text()
                    msg = json.loads(raw)
                    user_text = msg.get("message", "").strip()
                    raw_selected = msg.get("selected_skills", [])
                    if not user_text:
                        continue

                    selected = []
                    if isinstance(raw_selected, list):
                        selected = [str(x).strip() for x in raw_selected if str(x).strip()]
                    if not selected:
                        selected = _infer_skills_from_user_text(user_text, skills)

                    selected_ctx = _selected_skills_context(selected, skills)
                    user_payload = user_text
                    if selected_ctx:
                        user_payload = (
                            f"{user_text}\n\n"
                            "[Selected skill references from local SKILL.md files]\n"
                            f"{selected_ctx}"
                        )

                    session.history.append({"role": "user", "content": user_payload})
                    try:
                        reply = await _respond(websocket, session)
                        session.history.append({"role": "assistant", "content": reply})
                    except Exception as exc:
                        err_msg = f"[Agent error: {exc}]"
                        await websocket.send_text(
                            json.dumps({"type": "error", "message": str(exc)})
                        )
                        session.history.append({"role": "assistant", "content": err_msg})

            except WebSocketDisconnect:
                pass

        except Exception as e:
            print(f"[WS] Error occurred: {type(e).__name__}: {e}", flush=True)
            import traceback
            traceback.print_exc()

    return app


# ── Entry point ────────────────────────────────────────────────────────────────

def run_server(host: str = DEFAULT_HOST, port: int = DEFAULT_PORT) -> None:
    """Start the uvicorn server (blocking call — returns only when the server stops)."""
    _require_webdeps()
    import uvicorn  # type: ignore

    app = create_app()
    print(f"\n  NeuroClaw Web UI  →  http://{host}:{port}\n")
    uvicorn.run(app, host=host, port=port, log_level="info")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="NeuroClaw Web UI — start the browser-based chat interface."
    )
    parser.add_argument(
        "--host", default=DEFAULT_HOST,
        help=f"Bind host (default: {DEFAULT_HOST})",
    )
    parser.add_argument(
        "--port", type=int, default=DEFAULT_PORT,
        help=f"Port number (default: {DEFAULT_PORT})",
    )
    args = parser.parse_args()
    run_server(host=args.host, port=args.port)


if __name__ == "__main__":
    main()
